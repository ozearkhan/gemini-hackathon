"""Generate the PDLC Copilot hackathon deck as a real .pptx.

Reproducible: re-run to regenerate the deck after editing copy here. Diagrams
are drawn as native, editable PowerPoint vector shapes (no mermaid/Node needed),
matching the .mmd sources in ./diagrams. Legend is consistent across all
diagrams: gray = agent-executed phase, amber = human gate, teal = artifact/tool,
red = structural guardrail, blue = shared state, dashed/white = next-phase.

Honesty rule (Part 4 of the brief): rows/cells that are designed-but-not-wired
are explicitly labelled "next phase" — verified against the repo code, not the
brief's assumptions.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).with_name("PDLC-Copilot-Hackathon-Deck.pptx")

# ---- palette -------------------------------------------------------------
BRAND_DARK = RGBColor(0x14, 0x1E, 0x3C)   # navy header
BRAND_TEAL = RGBColor(0x2A, 0x6F, 0x6A)   # accent
EPAM_GREEN = RGBColor(0x6D, 0xBE, 0x45)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x20, 0x21, 0x24)
GRAY_TXT = RGBColor(0x5F, 0x63, 0x68)

KIND = {
    "agent":   (RGBColor(0xE8, 0xEA, 0xED), RGBColor(0x5F, 0x63, 0x68), INK),
    "gate":    (RGBColor(0xFC, 0xE8, 0xB2), RGBColor(0xB0, 0x60, 0x00), RGBColor(0x3D, 0x2C, 0x00)),
    "tool":    (RGBColor(0xA7, 0xD7, 0xD2), RGBColor(0x2A, 0x6F, 0x6A), RGBColor(0x0C, 0x30, 0x2D)),
    "artifact":(RGBColor(0xA7, 0xD7, 0xD2), RGBColor(0x2A, 0x6F, 0x6A), RGBColor(0x0C, 0x30, 0x2D)),
    "guard":   (RGBColor(0xF4, 0xC7, 0xC3), RGBColor(0xA5, 0x0E, 0x0E), RGBColor(0x3D, 0x0A, 0x0A)),
    "state":   (RGBColor(0xD7, 0xE3, 0xFC), RGBColor(0x1A, 0x56, 0xC4), RGBColor(0x0B, 0x2A, 0x66)),
    "planned": (WHITE, RGBColor(0x9A, 0xA0, 0xA6), RGBColor(0x3C, 0x40, 0x43)),
    "note":    (RGBColor(0xFE, 0xF7, 0xE0), RGBColor(0xB0, 0x60, 0x00), RGBColor(0x3D, 0x2C, 0x00)),
}

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]
_num = 0


def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _no_line(shape):
    shape.line.fill.background()


def _line(shape, rgb, width=1.0, dashed=False):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width)
    if dashed:
        ln = shape.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(d)


def _txt(shape, text, size=11, color=INK, bold=False, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        # sub-lines starting with '~' render italic/smaller (artifact captions)
        italic = ln.startswith("~")
        content = ln[1:] if italic else ln
        r = p.add_run(); r.text = content
        r.font.size = Pt(size - 2 if italic else size)
        r.font.bold = bold and not italic
        r.font.italic = italic
        r.font.color.rgb = GRAY_TXT if italic else color


def box(slide, x, y, w, h, text, kind="agent", size=11, bold=False,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    fill, line, txt = KIND[kind]
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(sp, fill)
    _line(sp, line, 1.25, dashed=(kind == "planned"))
    sp.shadow.inherit = False
    _txt(sp, text, size=size, color=txt, bold=bold)
    return sp


def conn(slide, x1, y1, x2, y2, color=GRAY_TXT, dashed=False, width=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return c


def free_text(slide, x, y, w, h, text, size=12, color=INK, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _txt(tb, text, size=size, color=color, bold=bold, align=align, anchor=anchor)
    return tb


def bullets(slide, x, y, w, h, items, size=14, gap_before=6):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        lvl, text = (it if isinstance(it, tuple) else (0, it))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap_before)
        bullet = "—  " if lvl == 0 else "•  "
        r = p.add_run(); r.text = bullet + text
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = INK
    return tb


def legend(slide, y, items):
    """items: list of (kind, label). Draw a compact swatch row."""
    x = 0.6
    for kind, label in items:
        fill, line, _ = KIND[kind]
        sw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(0.28), Inches(0.2))
        _fill(sw, fill); _line(sw, line, 1, dashed=(kind == "planned"))
        sw.shadow.inherit = False
        tb = slide.shapes.add_textbox(Inches(x + 0.32), Inches(y - 0.03), Inches(2.6), Inches(0.28))
        _txt(tb, label, size=9.5, color=GRAY_TXT, align=PP_ALIGN.LEFT)
        x += 0.42 + 0.03 * len(label) + 1.55


def chrome(slide, title=None):
    global _num
    _num += 1
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.5))
    _fill(band, BRAND_DARK); _no_line(band); band.shadow.inherit = False
    _txt(band, "GEMINI ENTERPRISE GLOBAL HACKATHON", size=11, color=WHITE,
         bold=True, align=PP_ALIGN.LEFT)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.5), prs.slide_width, Inches(0.045))
    _fill(accent, BRAND_TEAL); _no_line(accent); accent.shadow.inherit = False
    foot = slide.shapes.add_textbox(Inches(0.3), Inches(7.12), Inches(6), Inches(0.3))
    _txt(foot, "EPAM Proprietary & Confidential", size=9, color=GRAY_TXT, align=PP_ALIGN.LEFT)
    pg = slide.shapes.add_textbox(Inches(12.4), Inches(7.12), Inches(0.7), Inches(0.3))
    _txt(pg, str(_num), size=9, color=GRAY_TXT, align=PP_ALIGN.RIGHT)
    if title:
        t = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.3), Inches(0.6))
        _txt(t, title, size=24, color=BRAND_DARK, bold=True, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP)
    return slide


def new():
    return prs.slides.add_slide(BLANK)


# ===================================================================== SLIDE 1
s = new()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
_fill(band, BRAND_DARK); _no_line(band); band.shadow.inherit = False
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.05), prs.slide_width, Inches(0.06))
_fill(strip, EPAM_GREEN); _no_line(strip); strip.shadow.inherit = False
free_text(s, 0.8, 0.55, 12, 0.5, "GEMINI ENTERPRISE GLOBAL HACKATHON  ·  STREAM 2",
          size=13, color=EPAM_GREEN, bold=True)
free_text(s, 0.8, 1.35, 11.8, 1.4, "PDLC Copilot", size=54, color=WHITE, bold=True)
free_text(s, 0.8, 2.35, 11.8, 0.7,
          "A supervised, doc-grounded agent graph that owns the long tail of data-engineering requests",
          size=18, color=RGBColor(0xC5, 0xCB, 0xD8))
free_text(s, 0.8, 3.35, 11.8, 0.5, "Team  agenti-1711  —  Agentic PDLC", size=18,
          color=WHITE, bold=True)
free_text(s, 0.8, 3.95, 11.8, 1.6,
          "[NAME — Lead Architect]      [NAME — Agent / Tool Engineer]\n"
          "[NAME — Platform & Deploy]      [NAME — Eval & QA]",
          size=14, color=RGBColor(0xC5, 0xCB, 0xD8))
free_text(s, 0.8, 6.6, 12, 0.4, "EPAM Proprietary & Confidential", size=9,
          color=RGBColor(0x8A, 0x90, 0xA0))

# ===================================================================== SLIDE 2 — Problem
s = new(); chrome(s, "The Problem")
free_text(s, 0.6, 1.35, 12.1, 0.6,
          "Every \u201cquick\u201d data ask still needs the same six process steps a big platform build "
          "needs \u2014 but under deadline pressure, teams skip them. That is exactly where technical "
          "debt, compliance misses, and rework come from.", size=15, color=INK, bold=True)
bullets(s, 0.7, 2.7, 12.0, 3.6, [
    "It is a governance gap, not a tooling gap. The PDLC checklist already exists (our internal "
    "playbook) \u2014 what is missing is consistent execution at the volume of requests that actually arrive.",
    "The long tail is where the calendar time hides. A 10-ticker tracker, a competitor-watch dashboard, "
    "a one-off ingestion \u2014 individually small, collectively the bulk of a platform team's queue.",
    "Naive AI makes it worse. A generic copilot will confidently hallucinate a stock API's rate limit or "
    "invent Terraform syntax \u2014 turning a time problem into a trust problem.",
    "The steps that get skipped are the valuable ones: requirement sign-off, ADRs, review gates, "
    "JIRA traceability \u2014 not the code.",
], size=15)
free_text(s, 0.7, 6.35, 12, 0.4,
          "[Illustrative] plug in real org numbers if available: requests/quarter of this size, "
          "% with no written requirement doc, % with no ADR.", size=11, color=GRAY_TXT)

# ===================================================================== SLIDE 3 — Solution
s = new(); chrome(s, "The Solution")
free_text(s, 0.6, 1.3, 12.1, 0.6,
          "An agent graph that runs the full PDLC \u2014 and cannot skip a governance step even when "
          "a human would be tempted to.", size=15, color=INK, bold=True)
# mini pipeline strip
seq = [("Phase 0\nIntake", "agent"), ("Phase 1\nRequirements", "agent"),
       ("Gate-1\nsign-off", "gate"), ("Phase 2\nArchitecture", "agent"),
       ("Gate-2\nreview", "gate"), ("Phase 5\nInfra-as-Code", "agent"),
       ("Phase 4\nJIRA tree", "agent")]
x = 0.6
for i, (label, kind) in enumerate(seq):
    box(s, x, 2.25, 1.55, 0.7, label, kind, size=10, bold=True)
    if i < len(seq) - 1:
        conn(s, x + 1.55, 2.6, x + 1.72, 2.6, color=BRAND_TEAL)
    x += 1.72
bullets(s, 0.7, 3.35, 6.0, 3.4, [
    "How it helps:",
    (1, "Standardizes governance \u2014 nothing skips a gate."),
    (1, "Compresses elapsed time \u2014 parallel research + drafting vs. sequential meetings."),
    (1, "Artifacts (req doc, ADRs, JIRA tree) fall out as a byproduct of the work, not extra overhead."),
], size=14)
bullets(s, 6.9, 3.35, 5.9, 3.4, [
    "Who it helps:",
    (1, "Platform teams drowning in \u201cquick ask\u201d volume."),
    (1, "PMs/stakeholders who need a fast, defensible \u201ccan we build this & what will it cost.\u201d"),
    (1, "Junior engineers \u2014 a consistent example of \u201cgood\u201d every time."),
], size=14)
free_text(s, 0.7, 6.55, 12, 0.4,
          "Reference case worked end-to-end: \u201cDaily stock tracker for competitor watch.\u201d",
          size=13, color=BRAND_TEAL, bold=True)

# ===================================================================== SLIDE 4a — HLD
s = new(); chrome(s, "Architecture \u2014 HLD: System (what we shipped on Google ADK)")
legend(s, 1.28, [("agent", "agent phase"), ("gate", "human gate"),
                 ("tool", "tool / MCP"), ("guard", "guardrail (refuses)")])
box(s, 5.4, 1.65, 2.5, 0.62, "pdlc_coordinator · root_agent\nCoordinator-Dispatcher (fast model)", "agent", size=10, bold=True)
# phase row
px = 0.55; py = 2.7; pw = 2.3; ph = 0.75
phases = [("Phase 0\nintake_triage\n~triage note", "agent"),
          ("Phase 1\nrequirements_analyst\n~Requirement v1.x", "agent"),
          ("Phase 2\narchitecture\n~HLD+cost+ADRs", "agent"),
          ("Phase 5\niac (first)\n~Terraform in repo", "agent"),
          ("Phase 4\njira_planner\n~Epic→…→Task", "agent")]
xs = []
for i, (label, kind) in enumerate(phases):
    x = px + i * (pw + 0.2)
    xs.append(x)
    box(s, x, py, pw, ph, label, kind, size=9.5, bold=True)
    conn(s, 6.65, 2.27, x + pw / 2, py, color=GRAY_TXT, width=1)
# gates
box(s, xs[1] + pw + 0.2 - 2.3, 3.75, 2.3, 0.55, "Gate-1 · sign-off\nblocks until approval in state", "gate", size=9)
box(s, xs[2] + pw + 0.2 - 2.3, 3.75, 2.3, 0.55, "Gate-2 · review\napprove/conditional/reject", "gate", size=9)
# tool layer
tl = ["decide_load_pattern", "estimate_gcp_cost", "get_architecture_standard",
      "Dev-Knowledge MCP", "generate_terraform", "check_traceability", "researcher · google_search"]
tx = 0.55
for i, t in enumerate(tl):
    w = 1.72
    box(s, tx, 4.55, w, 0.55, t, "tool", size=9)
    tx += w + 0.05
# guardrails
box(s, 0.55, 5.35, 5.9, 0.5, "Structural guardrails (before_tool_callback, RAISES): approval gate · grounding gate · tool-call ceiling", "guard", size=9.5, bold=True)
box(s, 6.6, 5.35, 6.2, 0.5, "Observability: Cloud Trace (OTel) + BigQuery Agent Analytics \u2014 one nested trace per request_id", "state", size=9.5)
free_text(s, 0.55, 5.95, 12.3, 0.7,
          "Why centralized (not swarm): the PDLC is a paper trail \u2014 one coordinator, one place state "
          "and trace live, every choice auditable. Unhappy paths: gate reject loops back with the reason; "
          "missing approval/grounding \u2192 the persist tool refuses (fail-closed).",
          size=10.5, color=GRAY_TXT)

# ===================================================================== SLIDE 4b — LLD
s = new(); chrome(s, "Architecture \u2014 LLD: The Doc-Grounding Gate (the differentiator)")
legend(s, 1.28, [("agent", "agent step"), ("gate", "gate / decision"),
                 ("artifact", "artifact / state write"), ("guard", "refusal path")])
free_text(s, 0.6, 1.6, 12.1, 0.5,
          "This is a Python before_tool_callback that raises \u2014 a structural refusal, not a prompt. "
          "The persist step is architecturally unreachable without prior real research.",
          size=13, color=INK, bold=True)
row_y = 2.55
box(s, 0.55, row_y, 2.3, 0.85, "specialist plans\na factual claim /\ndesign decision", "agent", size=10)
box(s, 3.15, row_y, 2.4, 0.85, "researcher_agent\n(google_search) OR\nDev-Knowledge MCP fetch", "agent", size=10)
box(s, 5.85, row_y, 2.5, 0.85, "grounding gate\nenforce_research_grounded()\nresearcher consulted?", "gate", size=10, bold=True, shape=MSO_SHAPE.HEXAGON)
box(s, 8.7, row_y, 2.1, 0.85, "save_* tool\n(req / arch / terraform)", "agent", size=10)
box(s, 11.0, row_y, 1.85, 0.85, "published artifact\n~written to repo", "artifact", size=10)
conn(s, 2.85, row_y + 0.42, 3.15, row_y + 0.42, BRAND_TEAL)
conn(s, 5.55, row_y + 0.42, 5.85, row_y + 0.42, BRAND_TEAL)
conn(s, 8.35, row_y + 0.42, 8.7, row_y + 0.42, RGBColor(0x1E, 0x8E, 0x3E))
conn(s, 10.8, row_y + 0.42, 11.0, row_y + 0.42, BRAND_TEAL)
# refusal path
box(s, 5.85, 3.95, 2.5, 0.7, "GroundingRequiredError\n(+ ApprovalRequiredError if\nno recorded approval)", "guard", size=9.5)
conn(s, 7.1, row_y + 0.85, 7.1, 3.95, RGBColor(0xA5, 0x0E, 0x0E))
free_text(s, 8.5, 3.95, 4.3, 0.7, "absent \u2192 REFUSED,\nnothing persists (fail-closed)",
          size=11, color=RGBColor(0xA5, 0x0E, 0x0E), bold=True)
# state
box(s, 0.55, 5.0, 12.25, 1.05,
    "ADK session.state \u2014 single source of truth threaded through every node:\n"
    "grounding_recorded {researcher: true}   ·   approved_designs {slug: true}   ·   "
    "requirement_doc / adr_log / jira_tree   ·   _turn_tool_call_count", "state", size=11)
free_text(s, 0.55, 6.2, 12.2, 0.5,
          "\u201cOnly reachable path\u201d: the edge into save_* IS the gate. after_agent_callback on the "
          "researcher writes the evidence the gate checks \u2014 no evidence, no persist.",
          size=10.5, color=GRAY_TXT)

# ===================================================================== SLIDE 4c — Vision vs Build
s = new(); chrome(s, "Original Vision vs. Hackathon Build \u2014 same architecture, two runtimes")
free_text(s, 0.6, 1.3, 12.2, 0.5,
          "We prototyped the full vision on LangGraph/LangChain/MCP/LangSmith; for the timebox we "
          "re-platformed the same graph onto Google ADK / GCP. Honest status in the right column.",
          size=12.5, color=INK, bold=True)
rows = [
    ("Original design (LangChain)", "Role", "Hackathon build (Google ADK / GCP)", "Status"),
    ("LangGraph supervisor + subgraphs", "Orchestration / state machine", "Google ADK Coordinator-Dispatcher (LlmAgent + sub_agents)", "REAL"),
    ("LangChain tool layer", "Tool-calling glue", "ADK native tool-calling / FunctionTool", "REAL"),
    ("Context7 doc-gate", "Fresh grounding before any code/decision", "google_search researcher + Google Developer Knowledge MCP + repo standard tool", "REAL"),
    ("MCP servers (Atlassian, GitHub)", "Publish docs / tickets / PRs", "Artifacts emitted to repo; live SaaS MCP = sign-off gated", "NEXT PHASE"),
    ("LangSmith trace + grounding eval", "Observability + CI hallucination gate", "Cloud Trace + BQ Analytics wired; CI-blocking grounding eval", "TRACE REAL / EVAL NEXT"),
    ("interrupt() + Postgres checkpoint", "Durable pause/resume for sign-off", "before_tool_callback refusal in chat turn (durable gate = next)", "GATE REAL / DURABLE NEXT"),
    ("Typed PDLCState", "Single source of truth", "ADK session.state (shared dict)", "REAL"),
]
tbl = s.shapes.add_table(len(rows), 4, Inches(0.5), Inches(1.95),
                         Inches(12.33), Inches(4.6)).table
tbl.columns[0].width = Inches(3.3); tbl.columns[1].width = Inches(2.9)
tbl.columns[2].width = Inches(4.2); tbl.columns[3].width = Inches(1.93)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = val
        run.font.size = Pt(10.5 if r == 0 else 9.5)
        run.font.bold = (r == 0) or (c == 3)
        if r == 0:
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = BRAND_DARK
        else:
            status = row[3]
            if c == 3 and status == "REAL":
                cell.fill.solid(); cell.fill.fore_color.rgb = KIND["tool"][0]
                run.font.color.rgb = KIND["tool"][2]
            elif c == 3:
                cell.fill.solid(); cell.fill.fore_color.rgb = KIND["gate"][0]
                run.font.color.rgb = KIND["gate"][2]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xF2, 0xF4, 0xF7)
                run.font.color.rgb = INK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(5); cell.margin_right = Pt(5)
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
free_text(s, 0.5, 6.65, 12.3, 0.4,
          "Proof grounding is load-bearing, not decorative: the architecture agent cited the standard on a "
          "matching request and did NOT load it on a non-matching one \u2014 conditional, not always-on.",
          size=10.5, color=BRAND_TEAL, bold=True)

# ===================================================================== SLIDE 5 — Phase Flow
s = new(); chrome(s, "How We Planned It \u2014 PDLC Phase Flow (Phases 0\u20135, gates)")
legend(s, 1.28, [("agent", "agent-executed phase"), ("gate", "human gate"),
                 ("artifact", "published artifact")])
flow = [("Phase 0\nIntake & Triage", "agent", "triage note"),
        ("Phase 1\nRequirement", "agent", "Requirement v1.x"),
        ("Gate-1\nSign-off", "gate", ""),
        ("Phase 2\nArchitecture & HLD", "agent", "HLD + ADR log + cost"),
        ("Gate-2\nArch review", "gate", "review decision"),
        ("Phase 5\nInfra-as-Code", "agent", "Terraform in repo"),
        ("Phase 4\nJIRA breakdown", "agent", "Epic→…→Task + AC"),
        ("Phase 5b\nBuild→Review→Deploy", "agent", "PR + tests + runbook")]
fx = 0.55; fy = 2.2; fw = 1.5; fh = 0.85
for i, (label, kind, art) in enumerate(flow):
    box(s, fx, fy, fw, fh, label, kind, size=9, bold=True)
    if art:
        b = box(s, fx, fy + fh + 0.15, fw, 0.55, art, "artifact", size=8.5)
    if i < len(flow) - 1:
        conn(s, fx + fw, fy + fh / 2, fx + fw + 0.1, fy + fh / 2, color=BRAND_TEAL)
    fx += fw + 0.1
free_text(s, 0.55, 4.15, 12.3, 0.5,
          "Same backbone whether a human or an agent runs it \u2014 deliberate: it is what makes each phase "
          "independently automatable and auditable.", size=12.5, color=INK, bold=True)
bullets(s, 0.7, 4.75, 12.0, 1.9, [
    "Unhappy paths are first-class: Gate-1 revise loop (v1.0\u2192v1.1); Gate-2 conditional/reject loops "
    "back to Phase 2 with the reason recorded.",
    "Infra-first optimization: the agent provisions infra (Phase 5) right after Gate-2, before the JIRA "
    "breakdown \u2014 so there is no \u201cset up the infra\u201d ticket; the environment already exists.",
    "Every task/AC traces to a Phase-1 requirement or an ADR (check_task_traceability enforces it).",
], size=13)

# ===================================================================== SLIDE 6 — MCP Map
s = new(); chrome(s, "Grounding & MCP Ecosystem \u2014 wired vs. next-phase (honest)")
legend(s, 1.28, [("tool", "WIRED \u00b7 real in demo"), ("planned", "NEXT PHASE \u00b7 needs sign-off")])
box(s, 5.15, 1.75, 3.0, 0.6, "PDLC agents\n(architecture / iac / jira / requirements)", "agent", size=10, bold=True)
# wired
box(s, 0.55, 2.9, 3.9, 1.0, "Google Developer Knowledge MCP\nGoogle's own current docs \u00b7 Phase 2/4/5\nSTATIC \u2014 always available to codegen/ADR", "tool", size=10)
box(s, 4.7, 2.9, 3.9, 1.0, "google_search researcher sub-agent\nthird-party vendor/API facts \u00b7 Phase 1/2\nreal web grounding, not memory", "tool", size=10)
box(s, 8.85, 2.9, 3.95, 1.0, "get_architecture_standard tool\nrepo standard doc (Confluence stand-in)\nPhase 2 \u2014 deterministic read", "tool", size=10)
for cx in (2.5, 6.65, 10.8):
    conn(s, 6.65, 2.35, cx, 2.9, color=BRAND_TEAL)
# next phase
box(s, 0.55, 4.35, 3.9, 0.95, "Atlassian MCP\nJira tickets (P4) · Confluence publish (P1/P2)", "planned", size=10)
box(s, 4.7, 4.35, 3.9, 0.95, "GitHub MCP\nPR creation (P5)", "planned", size=10)
box(s, 8.85, 4.35, 3.95, 0.95, "Pattern-conditional MCPs\ndbt / Airflow / BigQuery \u2014 loaded only if\nPhase-2 pattern needs them", "planned", size=10)
for cx in (2.5, 6.65, 10.8):
    conn(s, 6.65, 2.35, cx, 4.35, color=RGBColor(0x9A, 0xA0, 0xA6), dashed=True, width=1)
box(s, 0.55, 5.55, 12.25, 0.85,
    "Today, the same paper trail with zero external dependency: Jira tree, requirement doc, HLD and "
    "Terraform are emitted as repo (GitLab) artifacts, not pushed to external SaaS. EPAM guardrail: "
    "external non-GCP SaaS MCP servers need Platform-Architect sign-off \u2014 so the default is GCP-native.",
    "note", size=11)

# ===================================================================== SLIDE 7 — Next Steps
s = new(); chrome(s, "Next Steps \u2014 closing the gaps in the vision\u2192build table")
bullets(s, 0.7, 1.45, 12.1, 5.2, [
    "1 · Stand up the grounding evaluator as a real CI-blocking gate \u2014 a run where a codegen step had "
    "no preceding doc-fetch fails the PR, exactly like a failing unit test (trace is wired; the gate is next).",
    "2 · Dynamic MCP binding \u2014 Phase 5 loads only the MCP servers the chosen architecture pattern "
    "actually needs (dbt / Airflow / BigQuery), not a fixed set.",
    "3 · Durable human gates \u2014 replace the in-session approval with an interrupt()-style durable pause "
    "(Firestore/Cloud SQL checkpoint) so a reviewer can approve async, days later, without losing state.",
    "4 · Enterprise-readiness checklist \u2014 Secret Manager for API keys, IAM scoping per MCP server, "
    "audit logging of every human-gate decision, and a defined owner/on-call model for the orchestrator itself.",
    "5 · Beyond Pattern D \u2014 prove the same graph selects and executes the heavier Patterns A\u2013C "
    "(Composer/Airflow, Dataflow, Dataproc) for larger requests, not just the lightweight long tail.",
    "6 · Live SaaS actions \u2014 opt-in Atlassian/GitHub MCP (Jira tickets, Confluence, PRs) once "
    "Platform-Architect sign-off lands \u2014 a one-time configured registration, thanks to the clean seams.",
], size=14, gap_before=10)

# ===================================================================== SLIDE 8 — Lessons
s = new(); chrome(s, "Lessons Learned")
bullets(s, 0.7, 1.45, 12.1, 5.2, [
    "The doc-gate thesis happened to us, not just in the demo: a locally-written TDD instruction "
    "contradicted google-agents-cli-eval's actual guidance \u2014 we fixed our instruction to match the "
    "vendor skill rather than keep our own. \u201cGrounding beats a hand-written instruction, even ours.\u201d",
    "The hardest part was not the LLM calls \u2014 it was making the human gates load-bearing rather than "
    "decorative under a timebox. Moving refusal into a before_tool_callback that raises was the turning point.",
    "Re-platforming LangGraph\u2192ADK clarified what was framework-specific vs. genuinely architectural: the "
    "doc-gate pattern and the typed state schema travelled cleanly; anything tied to LangSmith's evaluator API did not.",
    "Narrow, scope-limited specialists beat one monolithic agent \u2014 less surface area, measurably less "
    "hallucination, and each phase independently testable.",
    "[Add one more that is actually true for your team \u2014 do not fabricate specifics.]",
], size=15, gap_before=12)

# ===================================================================== SLIDE 9 — Visuals / Demo
s = new(); chrome(s, "Visuals / Demo")
ph = [("[SCREENSHOT: requirement doc generated \u2014 v1.0 with Open Questions]", 0.6, 1.5),
      ("[SCREENSHOT: human approval gate \u2014 record_human_approval in the chat turn]", 6.85, 1.5),
      ("[SCREENSHOT: grounding gate firing \u2014 GroundingRequiredError on save without research]", 0.6, 4.05),
      ("[SCREENSHOT: generated Terraform + JIRA tree written to the repo]", 6.85, 4.05)]
for label, x, y in ph:
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(5.9), Inches(2.35))
    _fill(b, RGBColor(0xF2, 0xF4, 0xF7)); _line(b, RGBColor(0x9A, 0xA0, 0xA6), 1, dashed=True)
    b.shadow.inherit = False
    _txt(b, label, size=12, color=GRAY_TXT, bold=True)
free_text(s, 0.6, 6.55, 12.2, 0.5,
          "Live demo prompts (all wired today): \u201cwhat's missing before we can build?\u201d  ·  "
          "\u201ccompare lightweight vs enterprise-Airflow cost\u201d  ·  \u201cscaffold the infra for the approved design.\u201d",
          size=11.5, color=BRAND_TEAL, bold=True)

prs.save(str(OUT))
print(f"Saved {OUT} with {_num + 1} slides")
